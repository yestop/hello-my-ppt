#!/usr/bin/env python3
"""生成图表参考文件 tests/projects/chart/reference/test-chart-all.pptx（python-pptx 官方库）。

覆盖 python-pptx 能生成的图表类型（每种一页）：
  bar(column) / line / area / pie / doughnut(≈官方 pie innerRadius) /
  scatter / bubble / radar

candlestick / waterfall / heatmap / treemap / sunburst / sankey 官方库不支持，
需手工拼 XML（见 HANDOFF C3）。

用途：C3 chart 官方化时解包比对 c:chart XML 结构。
"""
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE

prs = Presentation()
blank = prs.slide_layouts[6]

def add(ct, x, y, w, h, data, title=None):
    slide = prs.slides.add_slide(blank)
    gf = slide.shapes.add_chart(ct, x, y, w, h, data)
    chart = gf.chart
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    return slide

# 1. bar（簇状柱形）
cd = CategoryChartData()
cd.categories = ["Q1", "Q2", "Q3", "Q4"]
cd.add_series("收入", (120, 132, 101, 134))
cd.add_series("成本", (220, 182, 191, 234))
add(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(6), Inches(4), cd, "簇状柱形图")

# 2. line（折线 + 标记）
cd = CategoryChartData()
cd.categories = ["1月", "2月", "3月", "4月", "5月"]
cd.add_series("销量", (42, 55, 48, 68, 61))
add(XL_CHART_TYPE.LINE_MARKERS, Inches(1), Inches(1), Inches(6), Inches(4), cd, "折线图")

# 3. area（面积）
cd = CategoryChartData()
cd.categories = ["A", "B", "C", "D"]
cd.add_series("访问量", (30, 55, 42, 68))
add(XL_CHART_TYPE.AREA, Inches(1), Inches(1), Inches(6), Inches(4), cd, "面积图")

# 4. pie（饼图）
cd = CategoryChartData()
cd.categories = ["甲", "乙", "丙", "丁"]
cd.add_series("占比", (35, 25, 22, 18))
add(XL_CHART_TYPE.PIE, Inches(1), Inches(1), Inches(6), Inches(4), cd, "饼图")

# 5. doughnut（环形）
cd = CategoryChartData()
cd.categories = ["东", "西", "南", "北"]
cd.add_series("份额", (40, 30, 20, 10))
add(XL_CHART_TYPE.DOUGHNUT, Inches(1), Inches(1), Inches(6), Inches(4), cd, "环形图")

# 6. scatter（散点）
xd = XyChartData()
s1 = xd.add_series("样本A")
s1.add_data_point(1, 30)
s1.add_data_point(2, 55)
s1.add_data_point(3, 42)
s1.add_data_point(4, 68)
s1.add_data_point(5, 90)
s2 = xd.add_series("样本B")
s2.add_data_point(1.5, 20)
s2.add_data_point(2.5, 45)
s2.add_data_point(3.5, 60)
s2.add_data_point(4.5, 55)
add(XL_CHART_TYPE.XY_SCATTER, Inches(1), Inches(1), Inches(6), Inches(4), xd, "散点图")

# 7. bubble（气泡）
bd = BubbleChartData()
b1 = bd.add_series("气泡A")
b1.add_data_point(1, 30, 20)
b1.add_data_point(2, 55, 35)
b1.add_data_point(3, 42, 25)
b2 = bd.add_series("气泡B")
b2.add_data_point(1.5, 25, 30)
b2.add_data_point(2.5, 48, 40)
add(XL_CHART_TYPE.BUBBLE, Inches(1), Inches(1), Inches(6), Inches(4), bd, "气泡图")

# 8. radar（雷达）
cd = CategoryChartData()
cd.categories = ["速度", "力量", "防御", "敏捷", "耐力"]
cd.add_series("角色A", (80, 70, 60, 90, 75))
cd.add_series("角色B", (60, 85, 75, 55, 80))
add(XL_CHART_TYPE.RADAR, Inches(1), Inches(1), Inches(6), Inches(4), cd, "雷达图")

prs.save("tests/projects/chart/reference/test-chart-all.pptx")
print("已生成 tests/projects/chart/reference/test-chart-all.pptx（8 页）")
