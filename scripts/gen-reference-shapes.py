#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# gen-reference-shapes.py — 用 python-pptx 生成「全部预置形状」标准参考文件
# 输出：tests/projects/shape/reference/test-shapes-all.pptx（187 预置 + 自定义路径页）
# 布局与 tests/projects/shapes 项目一致（7 页 × 28 + 第 8 页自定义路径）
import re
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import _nsmap as nsmap

# 187 个预置名（从 preset-geometry.data.js 提取，保持与导出项目同源）
data = open('editor/core/preset-geometry.data.js', encoding='utf-8').read()
OUR_NAMES = re.findall(r'^  (\w+): \{', data, re.M)

# prst → MSO_SHAPE 成员
PRST2MSO = {m.xml_value: m for m in MSO_SHAPE.__members__.values() if getattr(m, 'xml_value', None)}
CONNECTOR_NAMES = ["line", "straightConnector1", "bentConnector2", "bentConnector3",
                   "bentConnector4", "bentConnector5", "curvedConnector2", "curvedConnector3",
                   "curvedConnector4", "curvedConnector5"]
missing_mso = [n for n in OUR_NAMES if n not in PRST2MSO and n not in CONNECTOR_NAMES]
assert not missing_mso, f"MSO_SHAPE 缺失: {missing_mso}"

prs = Presentation()
prs.slide_width = Pt(960)
prs.slide_height = Pt(540)
blank = prs.slide_layouts[6]
NS = 'xmlns:a="%s" xmlns:p="%s"' % (nsmap['a'], nsmap['p'])


def add_shape_xml(slide, sp_xml):
    el = parse_xml(sp_xml)
    slide.shapes._spTree.append(el)
    return el


def preset_sp(name, idx, x, y, w, h, fill="4472C4"):
    """构造 p:sp（无 p:style，与我们导出的简化结构一致——验证标注引线可见性）。"""
    off_x, off_y = int(x * 12700), int(y * 12700)
    ext_x, ext_y = int(w * 12700), int(h * 12700)
    return (
        f'<p:sp {NS}>'
        f'<p:nvSpPr><p:cNvPr id="{idx}" name="s{idx}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm><a:off x="{off_x}" y="{off_y}"/><a:ext cx="{ext_x}" cy="{ext_y}"/></a:xfrm>'
        f'<a:prstGeom prst="{name}"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>'
        f'</p:spPr>'
        f'</p:sp>'
    )


def custom_sp(idx, x, y, w, h, viewbox, path, fill="4472C4"):
    """自定义路径 custGeom（a:moveTo 全前缀写法）。"""
    off_x, off_y = int(x * 12700), int(y * 12700)
    ext_x, ext_y = int(w * 12700), int(h * 12700)
    vbw, vbh = viewbox
    return (
        f'<p:sp {NS}>'
        f'<p:nvSpPr><p:cNvPr id="{idx}" name="c{idx}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm><a:off x="{off_x}" y="{off_y}"/><a:ext cx="{ext_x}" cy="{ext_y}"/></a:xfrm>'
        f'<a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        f'<a:rect l="0" t="0" r="{vbw}" b="{vbh}"/>'
        f'<a:pathLst><a:path w="{vbw}" h="{vbh}">{path}</a:path></a:pathLst>'
        f'</a:custGeom>'
        f'<a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>'
        f'</p:spPr>'
        f'</p:sp>'
    )


# ---- 第 1~7 页：187 预置 ----
idx = 0
PER_PAGE = 28
for p in range(7):
    slide = prs.slides.add_slide(blank)
    chunk = OUR_NAMES[p * PER_PAGE:(p + 1) * PER_PAGE]
    for k, name in enumerate(chunk):
        idx += 1
        col, row = k % 5, k // 5
        x, y, w, h = 40 + col * 180, 74 + row * 112, 160, 92
        if name in CONNECTOR_NAMES:
            # 连接线类：XML 注入（MSO_SHAPE 枚举无连接线成员）
            add_shape_xml(slide, preset_sp(name, idx, x, y, w, h))
        else:
            sp = slide.shapes.add_shape(PRST2MSO[name], Pt(x), Pt(y), Pt(w), Pt(h))
            sp.name = f"s{idx}"
            sp.fill.solid()
            sp.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
            sp.line.fill.background()
    print(f"第 {p+1} 页: {len(chunk)} 个")

# ---- 第 8 页：自定义路径（与 shapes 项目同款）----
slide8 = prs.slides.add_slide(blank)
customs = [
    # (x, y, w, h, viewBox, path)
    (40, 74, 150, 150, (1000, 1000),
     '<a:moveTo><a:pt x="500" y="0"/></a:moveTo>'
     '<a:arcTo wR="500" hR="500" stAng="0" swAng="10800000"/>'
     '<a:arcTo wR="500" hR="500" stAng="10800000" swAng="10800000"/><a:close/>'
     '<a:moveTo><a:pt x="500" y="200"/></a:moveTo>'
     '<a:arcTo wR="300" hR="300" stAng="0" swAng="-10800000"/>'
     '<a:arcTo wR="300" hR="300" stAng="10800000" swAng="-10800000"/><a:close/>'),
    (220, 74, 240, 120, (1000, 500),
     '<a:moveTo><a:pt x="100" y="400"/></a:moveTo>'
     '<a:cubicBezTo><a:pt x="150" y="50"/><a:pt x="400" y="50"/><a:pt x="500" y="250"/></a:cubicBezTo>'
     '<a:cubicBezTo><a:pt x="600" y="450"/><a:pt x="800" y="400"/><a:pt x="900" y="100"/></a:cubicBezTo>'
     '<a:lnTo><a:pt x="900" y="400"/></a:lnTo><a:close/>'),
    (500, 74, 220, 130, (100, 60),
     '<a:moveTo><a:pt x="10" y="50"/></a:moveTo>'
     '<a:quadBezTo><a:pt x="50" y="10"/><a:pt x="90" y="50"/></a:quadBezTo>'
     '<a:quadBezTo><a:pt x="130" y="90"/><a:pt x="170" y="50"/></a:quadBezTo>'
     '<a:lnTo><a:pt x="210" y="50"/></a:lnTo>'
     '<a:lnTo><a:pt x="210" y="60"/></a:lnTo>'
     '<a:lnTo><a:pt x="170" y="60"/></a:lnTo>'
     '<a:quadBezTo><a:pt x="130" y="20"/><a:pt x="90" y="60"/></a:quadBezTo>'
     '<a:quadBezTo><a:pt x="50" y="100"/><a:pt x="10" y="60"/></a:quadBezTo><a:close/>'),
    (760, 74, 160, 120, (800, 600),
     '<a:moveTo><a:pt x="400" y="300"/></a:moveTo>'
     '<a:arcTo wR="300" hR="200" stAng="0" swAng="16200000"/><a:close/>'),
    (40, 250, 140, 110, (200, 160),
     '<a:moveTo><a:pt x="100" y="20"/></a:moveTo>'
     '<a:lnTo><a:pt x="180" y="80"/></a:lnTo>'
     '<a:lnTo><a:pt x="100" y="140"/></a:lnTo>'
     '<a:lnTo><a:pt x="20" y="80"/></a:lnTo><a:close/>'
     '<a:moveTo><a:pt x="100" y="50"/></a:moveTo>'
     '<a:lnTo><a:pt x="60" y="80"/></a:lnTo>'
     '<a:lnTo><a:pt x="100" y="110"/></a:lnTo>'
     '<a:lnTo><a:pt x="140" y="80"/></a:lnTo><a:close/>'),
    (220, 250, 180, 110, (360, 220),
     '<a:moveTo><a:pt x="0" y="110"/></a:moveTo>'
     '<a:arcTo wR="180" hR="110" stAng="0" swAng="10800000"/>'
     '<a:arcTo wR="180" hR="110" stAng="10800000" swAng="10800000"/><a:close/>'),
    (440, 250, 160, 110, (320, 220),
     '<a:moveTo><a:pt x="160" y="0"/></a:moveTo>'
     '<a:lnTo><a:pt x="320" y="220"/></a:lnTo>'
     '<a:lnTo><a:pt x="160" y="170"/></a:lnTo>'
     '<a:lnTo><a:pt x="0" y="220"/></a:lnTo><a:close/>'),
]
for i, (x, y, w, h, vb, path) in enumerate(customs):
    add_shape_xml(slide8, custom_sp(200 + i, x, y, w, h, vb, path))
print("第 8 页: 7 个自定义路径")

out = "tests/projects/shape/reference/test-shapes-all.pptx"
prs.save(out)
print(f"✓ 已生成 {out}（{len(OUR_NAMES)} 预置 + 7 自定义，8 页）")
